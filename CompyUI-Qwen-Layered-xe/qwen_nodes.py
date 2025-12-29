import torch
import gc
import numpy as np
import os
import tempfile
from PIL import Image
from pptx import Presentation
from diffusers import QwenImageLayeredPipeline
import folder_paths
from datetime import datetime

def cleanup():
    gc.collect()
    if torch.cuda.is_available():
        torch.cuda.empty_cache()

class QwenLayeredLoader:
    @classmethod
    def INPUT_TYPES(s):
        return {
            "required": {
                "model_name": (["Qwen/Qwen-Image-Layered"],),
                "low_vram": ("BOOLEAN", {"default": True}),
            }
        }
    
    RETURN_TYPES = ("QWEN_PIPE",)
    FUNCTION = "load_model"
    CATEGORY = "QwenLayered"

    def load_model(self, model_name, low_vram):
        cleanup()
        
        # モデルをロード
        pipeline = QwenImageLayeredPipeline.from_pretrained(
            model_name,
            torch_dtype=torch.bfloat16,
            device_map=None,
        )
        
        # 最適化設定
        pipeline.enable_attention_slicing()
        
        # 【重要】QwenのVAE(Conv3D)はCPUで不安定なため、tiling/slicingのみ有効化し、
        # デバイス移動はDiffusersのオフロード管理に完全に任せる
        pipeline.vae.enable_slicing()
        pipeline.vae.enable_tiling()

        if low_vram:
            # これにより、VAE実行時のみGPUに乗り、終わると自動でCPUに退避されます
            pipeline.enable_sequential_cpu_offload()
        else:
            pipeline.to("cuda")
            
        return (pipeline,)

class QwenLayeredInference:
    @classmethod
    def INPUT_TYPES(s):
        return {
            "required": {
                "pipeline": ("QWEN_PIPE",),
                "image": ("IMAGE",),
                "prompt": ("STRING", {"multiline": True, "default": ""}),
                "seed": ("INT", {"default": 0, "min": 0, "max": 0xffffffffffffffff}),
                "steps": ("INT", {"default": 50, "min": 1, "max": 100}),
                "guidance_scale": ("FLOAT", {"default": 4.0, "min": 1.0, "max": 10.0}),
                "layers": ("INT", {"default": 4, "min": 2, "max": 10}),
            }
        }

    RETURN_TYPES = ("IMAGE",)
    FUNCTION = "process"
    CATEGORY = "QwenLayered"

    def process(self, pipeline, image, prompt, seed, steps, guidance_scale, layers):
        cleanup()
        
        # ComfyUI Tensor [B,H,W,C] -> PIL (RGBA)
        i = 255. * image[0].cpu().numpy()
        pil_img = Image.fromarray(np.clip(i, 0, 255).astype(np.uint8)).convert("RGBA")

        # Generatorのデバイス設定
        # sequential_cpu_offload使用時は、モデルのデバイスを直接参照せず、
        # CPUまたは現在のメインデバイスを指定するのが安全です
        generator = torch.Generator(device="cpu").manual_seed(seed)

        with torch.inference_mode():
            output = pipeline(
                image=pil_img,
                generator=generator,
                true_cfg_scale=guidance_scale,
                prompt=prompt if prompt != "" else None,
                num_inference_steps=steps,
                layers=layers,
                resolution=640,
                cfg_normalize=True,
            )
            output_images = output.images[0]

        # Tensor変換
        output_tensors = []
        for img in output_images:
            img_np = np.array(img.convert("RGB")).astype(np.float32) / 255.0
            output_tensors.append(torch.from_numpy(img_np))
        
        cleanup()
        return (torch.stack(output_tensors),)

class QwenPPTXExport:
    @classmethod
    def INPUT_TYPES(s):
        return {
            "required": {
                "images": ("IMAGE",),
                "filename_prefix": ("STRING", {"default": "qwen_layers"}),
            }
        }

    RETURN_TYPES = ("STRING",)
    OUTPUT_NODE = True
    FUNCTION = "export_pptx"
    CATEGORY = "QwenLayered"

    def export_pptx(self, images, filename_prefix):
        # 現在の日時を YYYYMMDD_HHMMSS 形式で取得
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # ファイル名を「接頭辞_日時.pptx」にする
        filename = f"{filename_prefix}_{timestamp}.pptx"
        
        prs = Presentation()
        h, w = images[0].shape[0], images[0].shape[1]
        prs.slide_width = int((w / 96) * 914400)
        prs.slide_height = int((h / 96) * 914400)
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for i in range(len(images)):
            img_np = (images[i].cpu().numpy() * 255).astype(np.uint8)
            pil_img = Image.fromarray(img_np)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                pil_img.save(tmp.name)
                slide.shapes.add_picture(tmp.name, 0, 0, width=prs.slide_width, height=prs.slide_height)

        output_path = os.path.join(folder_paths.get_output_directory(), filename)
        prs.save(output_path)
        
        print(f"PPTX saved to: {output_path}") # コンソールに保存先を表示
        return (output_path,)